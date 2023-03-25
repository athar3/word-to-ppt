import docx
from pptx import Presentation
from pptx.util import Inches

# membaca file word
doc = docx.Document('filtes.docx')

# membuat presentasi powerpoint
prs = Presentation()

# loop untuk setiap paragraf pada file word
for para in doc.paragraphs:
    # menambahkan slide baru pada presentasi powerpoint
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    
    # menambahkan teks paragraf pada kotak teks di bawah judul pada slide
    
    text_frame = slide.placeholders[1].text_frame
    text_frame.text = ''
    for run in para.runs:
        text_frame.text += run.text
    
# menyimpan presentasi powerpoint
prs.save('filtes.pptx')
