from app import Flask, request, render_template
import os
import docx
from pptx import Presentation
from pptx.util import Inches

app = Flask(__name__)

@app.route('/')
def index():
    return render_template('D:\simpleprojek\wordtoppt\web\home.html')

@app.route('/convert', methods=['POST'])
def convert():
    file = request.files['file']
    file.save(file.filename)

    # mengubah file word menjadi ppt
    doc = docx.Document(file.filename)
    prs = Presentation()

    for para in doc.paragraphs:
        if len(para.text.strip()) == 1:
            slide = prs.slides.add_slide(prs.slide_layouts[0])
            slide.shapes.title.text = para.text.strip()
        else:
            slide = prs.slides.add_slide(prs.slide_layouts[1])
            text_frame = slide.shapes.placeholders[1].text_frame
            text_frame.text = para.text.strip()

    # menyimpan ppt
    new_filename = os.path.splitext(file.filename)[0] + '.pptx'
    prs.save(new_filename)

    return f'The file {new_filename} has been created.'

if __name__ == '__main__':
    app.run(debug=True)
