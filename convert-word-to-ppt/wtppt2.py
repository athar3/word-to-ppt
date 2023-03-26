import docx
from pptx import Presentation
from pptx.util import Inches

# membaca file word
doc = docx.Document('uas.docx')

# membuat presentasi powerpoint
prs = Presentation()

# loop untuk setiap paragraf pada file word
for para in doc.paragraphs:
    # menentukan apakah paragraf hanya terdiri dari satu baris atau tidak
    is_one_line = len(para.text.split('\n')) == 1
    
    # menambahkan slide baru pada presentasi powerpoint
    if is_one_line:
        slide = prs.slides.add_slide(prs.slide_layouts[0])
    else:
        slide = prs.slides.add_slide(prs.slide_layouts[1])
    
    # menambahkan teks paragraf pada kotak judul atau kotak teks di bawah judul pada slide
    if is_one_line:
        title_frame = slide.shapes.title.text_frame
        title_frame.text = para.text
    else:
        text_frame = slide.placeholders[1].text_frame
        text_frame.text = ''
        for run in para.runs:
            lines = run.text.split('\n')
            for i, line in enumerate(lines):
                if i == 0:
                    text_frame.text += line
                else:
                    p = text_frame.add_paragraph()
                    p.text = line
    
# menyimpan presentasi powerpoint
prs.save('tes.pptx')
