from flask import Flask, render_template, request, send_file
from docx import Document
from pptx import Presentation
from pptx.util import Inches
import os

app = Flask(__name__)

@app.route('/')
def home():
    return render_template('home.html')

@app.route('/convert', methods=['POST'])
def convert():
    file = request.files['file']
    filename = file.filename
    
    # cek apakah file yang diupload adalah file docx
    if not filename.lower().endswith('.docx'):
        return "File yang diupload harus berupa file docx"
    
    # simpan file yang diupload
    file.save(filename)
    
    # buka file docx
    doc = Document(filename)
    
    # buat presentasi powerpoint baru
    prs = Presentation()
    
    # loop untuk setiap paragraf pada file word
    for para in doc.paragraphs:
        # menambahkan slide baru pada presentasi powerpoint
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        
        # menambahkan teks paragraf pada kotak teks di bawah judul pada slide
        text_frame = slide.placeholders[1].text_frame
        text_frame.text = ''
        for run in para.runs:
            text_frame.text += run.text
    
    # simpan presentasi powerpoint dengan nama yang sama dengan file word
    pptx_filename = os.path.splitext(filename)[0] + '.pptx'
    prs.save(pptx_filename)
    
    # kirimkan file pptx yang sudah dihasilkan ke user untuk diunduh
    return send_file(pptx_filename, as_attachment=True)

if __name__ == '__main__':
    app.run(debug=True)
