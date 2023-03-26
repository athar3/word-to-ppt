import docx
from pptx import Presentation
from pptx.util import Inches

#baca file word
doc = docx.Document('filtes.docx')

#presentasi powerpoint
prs = Presentation()

# loop tiap paragraf di file word
for para in doc.paragraphs:
    # add slide baru di ppt
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    
    # add teks paragraf di kotak teks
    text_frame = slide.placeholders[1].text_frame
    text_frame.text = ''
    for run in para.runs:
        text_frame.text += run.text
    
# simpan ppt
prs.save('filtes.pptx')
